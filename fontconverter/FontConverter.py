from .convert_util import convert, fix_punc
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from aspose.slides import Presentation as slidesPresentation
from aspose.slides import util, export, FontData
import re


class FontConverter:
    def convert_font(self, text_frame, src, dest):
        if self.is_arabic(text_frame.text):
            return
        prgs = text_frame.paragraphs
        for prg in prgs:
            runs = prg.runs
            for i in range(len(runs)):
                if runs[i].font.name == src:
                    prev = ""
                    next = ""
                    if i > 0:
                        prev = runs[i - 1].text
                    if i < len(runs) - 1:
                        next = runs[i + 1].text
                    converted = convert(runs[i].text, src, dest)
                    prev, curr, next = fix_punc(converted, prev, next, src, dest)
                    if i > 0:
                        runs[i - 1].text = prev
                    if i < len(runs) - 1:
                        runs[i + 1].text = next
                    runs[i].text = curr

    def convert_all_text(self, file, src, dest):
        if src == "Unicode":
            src = "Coptic New Athanasius"
        if dest == "Unicode":
            dest = "Coptic New Athanasius"
        ppt = Presentation(file)
        for idx, slide in enumerate(ppt.slides):
            shapes = slide.shapes
            for shape in shapes:
                if shape.has_table:
                    for cell in shape.table.iter_cells():
                        self.convert_font(cell.text_frame, src, dest)
                if shape.has_text_frame:
                    self.convert_font(shape.text_frame, src, dest)
        ppt.save(file)

    def change_font(self, file, src, dest):
        if src == "Unicode":
            src = "Coptic New Athanasius"
        if dest == "Unicode":
            dest = "Coptic New Athanasius"
        ppt = slidesPresentation(file)
        textFramesPPTX = util.SlideUtil.get_all_text_frames(ppt, True)
        for idx, text_frame in enumerate(textFramesPPTX):
            if self.is_arabic(text_frame.text):
                continue
            prgs = text_frame.paragraphs
            for prg in prgs:
                for port in prg.portions:
                    if (
                        port.portion_format.latin_font is not None
                        and port.portion_format.latin_font.font_name == src
                    ):
                        if dest == "Coptic New Athanasius":
                            port.portion_format.latin_font = FontData(dest)
                            port.portion_format.east_asian_font = FontData(dest)
                        else:
                            port.portion_format.latin_font = FontData(dest)
                            port.portion_format.east_asian_font = None
        ppt.save(file, export.SaveFormat.PPTX)

    def is_arabic(self, text):
        pattern = re.compile(".*[\\u0600-\\u06FF]")
        match = pattern.match(text)
        return match is not None
