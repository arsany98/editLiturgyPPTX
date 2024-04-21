from pptx import Presentation


class TextReplacer:
    def replace_text(self, text_frame, find, replace):
        prgs = text_frame.paragraphs
        for prg in prgs:
            for run in prg.runs:
                run.text = run.text.replace(find, replace)

    def edit_ppt(self, file, find, replace):
        try:
            ppt = Presentation(file)
            for idx, slide in enumerate(ppt.slides):
                shapes = slide.shapes
                for shape in shapes:
                    if shape.has_table:
                        for cell in shape.table.iter_cells():
                            self.replace_text(cell.text_frame, find, replace)
                    elif shape.has_text_frame:
                        self.replace_text(shape.text_frame, find, replace)

            ppt.save(file)
        except Exception as e:
            print(file, "is invalid", e)
