from pptx.util import Cm, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE, MSO_COLOR_TYPE
from pptx import Presentation
import re


def CmOrNone(val):
    if val == None:
        return None
    return Cm(val)


def PtOrNone(val):
    if val == None:
        return None
    return Pt(val)


class PythonPPTXManager:
    def __init__(
        self,
        new_width,
        new_height,
        font_size_increase,
        copt_font_size_increase,
        exclude_first_slide,
        exclude_outlined,
        shape_margin,
        table_margin,
        line_width,
        textbox_position,
        extend_textbox_width
    ):
        self.new_width = CmOrNone(new_width)
        self.new_height = CmOrNone(new_height)
        self.font_size_increase = PtOrNone(font_size_increase)
        self.copt_font_size_increase = PtOrNone(copt_font_size_increase)
        self.line_width = PtOrNone(line_width)
        self.exclude_first_slide = exclude_first_slide
        self.exclude_outlined = exclude_outlined
        self.shape_margin = shape_margin
        self.table_margin = table_margin
        self.textbox_position = CmOrNone(textbox_position)
        self.extend_textbox_width = extend_textbox_width

    def change_shape_width(self, shape, exclude):
        ratio = self.new_width / self.old_width
        diff = self.new_width - self.old_width
        w = shape.width
        h = shape.height
        l = shape.left
        t = shape.top
        if exclude:
            shape.left = Emu(l + (diff / 2))
            return
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.left = Emu((l * ratio) if l > 0 else l)
        elif shape.is_placeholder:
            shape.width = Emu(w * ratio)
            shape.left = Emu((l * ratio) if l > 0 else l)
            shape.height = Emu(h)
            shape.top = Emu(t)
        else:
            shape.width = Emu(w * ratio)
            shape.left = Emu((l * ratio) if l > 0 else l)
        if shape.has_table:
            for column in shape.table.columns:
                column.width = Emu(column.width * ratio)

    def change_shape_height(self, shape, exclude):
        ratio = self.new_height / self.old_height
        diff = self.new_height - self.old_height
        w = shape.width
        h = shape.height
        l = shape.left
        t = shape.top
        if exclude:
            shape.top = Emu(t + (diff / 2))
            return
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.top = Emu((t * ratio) if t > 0 else t)
        elif shape.is_placeholder:
            shape.height = Emu(h * ratio)
            shape.top = Emu((t * ratio) if t > 0 else t)
            shape.width = Emu(w)
            shape.left = Emu(l)
        else:
            shape.height = Emu(h * ratio)
            shape.top = Emu((t * ratio) if t > 0 else t)
        if shape.has_table:
            for row in shape.table.rows:
                row.height = Emu(row.height * ratio)

    def increase_font_size(self, text_frame, font_increase):
        prgs = text_frame.paragraphs
        for prg in prgs:
            for run in prg.runs:
                if run.font.size is not None:  # error font less than 0
                    run.font.size += font_increase

    def increase_arabic_shape_font_size(self, shape):
        if shape.has_table:
            for cell in shape.table.iter_cells():
                if self.is_arabic(cell.text_frame.text):
                    self.increase_font_size(cell.text_frame, self.font_size_increase)
        if shape.has_text_frame and self.font_size_increase is not None:
            if self.is_arabic(shape.text_frame.text):
                self.increase_font_size(shape.text_frame, self.font_size_increase)

    def increase_coptic_shape_font_size(self, shape):
        if shape.has_table:
            for cell in shape.table.iter_cells():
                if not self.is_arabic(cell.text_frame.text):
                    self.increase_font_size(
                        cell.text_frame, self.copt_font_size_increase
                    )
        if shape.has_text_frame and self.font_size_increase is not None:
            if not self.is_arabic(shape.text_frame.text):
                self.increase_font_size(shape.text_frame, self.copt_font_size_increase)

    def edit_slide(self, slide, exclude=False):
        shapes = slide.shapes
        for shape in shapes:
            if self.new_width is not None:
                self.change_shape_width(shape, exclude)
            if self.new_height is not None:
                self.change_shape_height(shape, exclude)
            if not exclude and not (
                self.exclude_outlined
                and (
                    shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                    or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                )
                and shape.line.fill.type == MSO_FILL_TYPE.SOLID
            ):
                if self.font_size_increase is not None and self.font_size_increase != 0:
                    self.increase_arabic_shape_font_size(shape)
                if (
                    self.copt_font_size_increase is not None
                    and self.copt_font_size_increase != 0
                ):
                    self.increase_coptic_shape_font_size(shape)
            self.set_shape_margin(shape)
            self.set_line_width(shape)

    def edit_ppt(self, file):
        ppt = Presentation(file)
        self.old_width = ppt.slide_width
        if self.new_width is not None:
            ppt.slide_width = self.new_width
        self.old_height = ppt.slide_height
        if self.new_height is not None:
            ppt.slide_height = self.new_height

        self.edit_slide(ppt.slide_master)
        for slide_layout in ppt.slide_master.slide_layouts:
            self.edit_slide(slide_layout)

        for idx, slide in enumerate(ppt.slides):
            self.edit_slide(slide, exclude=(self.exclude_first_slide and idx == 0))
            self.move_textbox_to_position(slide)
            self.extend_textbox_width_to_match_slide(slide)

        ppt.save(file)

    def is_arabic(self, text):
        pattern = re.compile(".*[\\u0600-\\u06FF]")
        match = pattern.match(text)
        return match is not None

    def set_text_margin(self, text_frame, margin):
        if margin.left is not None:
            text_frame.margin_left = Cm(margin.left)
        if margin.top is not None:
            text_frame.margin_top = Cm(margin.top)
        if margin.right is not None:
            text_frame.margin_right = Cm(margin.right)
        if margin.bottom is not None:
            text_frame.margin_bottom = Cm(margin.bottom)

    def set_shape_margin(self, shape):
        if shape.has_table:
            for cell in shape.table.iter_cells():
                self.set_text_margin(cell, self.table_margin)
        if shape.has_text_frame:
            self.set_text_margin(shape.text_frame, self.shape_margin)

    def set_line_width(self, shape):
        if (
            self.line_width is not None
            and (
                shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            )
            and shape.line.fill.type == MSO_FILL_TYPE.SOLID
        ):
            shape.line.width = self.line_width

    def get_slides_count(file):
        ppt = Presentation(file)

        return len(ppt.slides)

    def apply_font(paragraph, font):
        p = paragraph
        p.font.size = font.size
        if font.color.type == MSO_COLOR_TYPE.RGB:
            p.font.color.rgb = font.color.rgb
        elif font.color.type == MSO_COLOR_TYPE.SCHEME:
            p.font.color.theme_color = font.color.theme_color
        p.font.bold = font.bold
        p.font.name = font.name

    def split_on_newline(file):
        ppt = Presentation(file)
        textboxes = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                    or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                ):
                    textboxes.append(shape)
                    break
        for i in range(len(textboxes) - 1):
            if textboxes[i].text_frame.text == textboxes[i + 1].text_frame.text:
                font = textboxes[i].text_frame.paragraphs[0].runs[0].font
                t = (
                    textboxes[i]
                    .text_frame.text.replace("\x0b\n", "\x0b\x0b")
                    .split("\x0b\x0b")
                )
                if len(t) >= 2:
                    if t[0].strip() != "":
                        textboxes[i].text_frame.clear()
                        textboxes[i].text_frame.paragraphs[0].text = t[0] + "\x0b"
                        PythonPPTXManager.apply_font(
                            textboxes[i].text_frame.paragraphs[0], font
                        )
                    if t[1].strip() != "":
                        textboxes[i + 1].text_frame.clear()
                        textboxes[i + 1].text_frame.paragraphs[0].text = t[1]
                        PythonPPTXManager.apply_font(
                            textboxes[i + 1].text_frame.paragraphs[0], font
                        )

        ppt.save(file)

    def move_textbox_to_position(self, slide):
        if self.textbox_position is None:
            return
        textboxes = 0
        for shape in slide.shapes:
            if (
                shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            ):
                textboxes += 1
        if textboxes == 1:
            for shape in slide.shapes:
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                    or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                ):
                    shape.top = self.textbox_position

    def extend_textbox_width_to_match_slide(self, slide):
        if self.extend_textbox_width is None:
            return
        textboxes = 0
        for shape in slide.shapes:
            if (
                shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            ):
                textboxes += 1
        if textboxes == 1:
            for shape in slide.shapes:
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                    or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                ):
                    shape.width = self.new_width if self.new_width else self.old_width
                    shape.left = 0