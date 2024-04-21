from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Emu


class TextBoxesFormatter:
    def __init__(self, line_width):
        self.line_width = line_width

    def edit_ppt(self, file):
        try:
            ppt = Presentation(file)
            for idx, slide in enumerate(ppt.slides):
                for shape in slide.shapes:
                    if (
                        shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                        or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                    ):
                        if self.line_width != 0:
                            shape.line.width = Pt(self.line_width)

            ppt.save(file)
        except Exception as e:
            print(file, "is invalid", e)
