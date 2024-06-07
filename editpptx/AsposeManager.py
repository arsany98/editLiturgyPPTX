from aspose.slides import (
    FontsLoader,
    export,
    FontData,
    Table,
    util,
    ShapesAlignmentType,
)
from aspose.slides import Presentation as slidesPresentation
from fontTools import ttLib
import os
from pptx import Presentation
from pptx.util import Cm


class AsposeManager:
    def remove_water_mark(self, file):
        ppt = Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if (
                    shape.has_text_frame
                    and shape.text_frame.text.lower().find("aspose") != -1
                ):
                    slide.shapes.element.remove(shape.element)
        ppt.save(file)

    def get_font_files(self, dir_path):
        result = []
        for root, dirs, files in os.walk(dir_path):
            for f in files:
                file = os.path.join(root, f)

                if file.endswith(".ttf"):
                    font = ttLib.TTFont(file)
                    result.append(FontData(font["name"].getDebugName(1)))
        return result

    def get_embedded_fonts(self, files):
        fonts = {}
        for file in files:
            presentation = slidesPresentation(file)
            embeddedFonts = presentation.fonts_manager.get_embedded_fonts()
            embeddedFonts = [(x.font_name, x) for x in embeddedFonts]
            fonts.update(embeddedFonts)
        return fonts

    def remove_embedded_font(self, file, font_data):
        try:
            presentation = slidesPresentation(file)
            presentation.fonts_manager.remove_embedded_font(font_data)
            presentation.save(file, export.SaveFormat.PPTX)
        except Exception as e:
            print(e)

    def embed_fonts(self, file, font_dir):
        presentation = slidesPresentation(file)
        allFonts = self.get_font_files(font_dir)
        embeddedFonts = presentation.fonts_manager.get_embedded_fonts()
        embeddedFonts = [x.font_name for x in embeddedFonts]

        FontsLoader.load_external_fonts([font_dir])
        for font in allFonts:
            if font.font_name not in embeddedFonts:
                try:
                    presentation.fonts_manager.add_embedded_font(
                        font, export.EmbedFontCharacters.ALL
                    )
                except Exception as e:
                    print(f"{font} {e}")
            else:
                print(f"{font} Font is already embedded.")
        presentation.save(file, export.SaveFormat.PPTX)

    def move_table_to_position(self, presentation, position):
        for slide in presentation.slides:
            tables = 0
            for shape in slide.shapes:
                if type(shape) == Table:
                    tables += 1
            for shape in slide.shapes:
                if tables == 1 and type(shape) == Table:
                    shape.y = Cm(position).pt - shape.rows[0].height

    def split_slides(self, presentation):
        slides_numbers = []
        for slide in presentation.slides:
            tables = 0
            for shape in slide.shapes:
                if type(shape) == Table:
                    tables += 1
            if tables == 2:
                slides_numbers.append(slide.slide_number)
        added_slides = 0
        for i in slides_numbers:
            presentation.slides.insert_clone(
                i + added_slides, presentation.slides[i + added_slides - 1]
            )
            added_slides += 1
            slide1 = presentation.slides[i + added_slides - 2]
            slide2 = presentation.slides[i + added_slides - 1]
            tables_indices = []
            for idx, shape in enumerate(slide1.shapes):
                if type(shape) == Table:
                    tables_indices.append(idx)

            slide1.shapes.remove_at(tables_indices[1])
            slide2.shapes.remove_at(tables_indices[0])

    def center_tables(self, presentation):
        for slide in presentation.slides:
            tables = 0
            for shape in slide.shapes:
                if type(shape) == Table:
                    tables += 1
            if tables == 1:
                for idx, shape in enumerate(slide.shapes):
                    if type(shape) == Table:
                        for row in shape.rows:
                            row.height
                        util.SlideUtil.align_shapes(
                            ShapesAlignmentType.ALIGN_MIDDLE, True, slide, [idx]
                        )

    def edit_ppt(self, file, table_position, split_slides, center_tables):
        if not table_position and not split_slides and not center_tables:
            return
        presentation = slidesPresentation(file)
        if split_slides:
            self.split_slides(presentation)
        if center_tables:
            self.center_tables(presentation)
        elif table_position is not None:
            self.move_table_to_position(presentation, table_position)
        presentation.save(file, export.SaveFormat.PPTX)
        self.remove_water_mark(file)
